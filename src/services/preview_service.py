from __future__ import annotations

import hashlib
import io
import platform
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol

from src.config import AppConfig


PREVIEW_RENDER_WIDTH = 1600


@dataclass(frozen=True)
class PreviewResult:
    kind: str
    content: bytes | str
    total_pages: int
    renderer: str
    page_no: int
    warning: str = ""


class PreviewRenderer(Protocol):
    name: str

    def can_render(self, file_path: Path) -> bool:
        ...

    def render(self, file_path: Path, page_no: int, token: str, matched_text: str = "") -> PreviewResult:
        ...


class PreviewService:
    def __init__(self, config: AppConfig) -> None:
        self.config = config
        self.pdf_renderer = PdfPageRenderer()
        self.preconverted_renderer = PreconvertedPreviewRenderer(self.pdf_renderer)
        self.powerpoint_renderer = PowerPointComRenderer()
        self.pptx_text_renderer = PptxTextRenderer()

    def file_cache_token(self, file_path: str) -> str:
        stat = Path(file_path).stat()
        return f"{stat.st_mtime_ns}:{stat.st_size}"

    def render(self, file_path: str, page_no: int, matched_text: str = "") -> PreviewResult:
        path = Path(file_path)
        token = self.file_cache_token(str(path))
        page_no = request_page_no(page_no)
        suffix = path.suffix.lower()
        if suffix == ".pdf":
            return self.pdf_renderer.render(path, page_no, token, matched_text)
        if suffix in {".ppt", ".pptx"}:
            return self._render_presentation(path, page_no, token, matched_text)
        return PreviewResult("text", matched_text, 1, "text", page_no)

    def prefetch(self, file_path: str, page_no: int) -> None:
        path = Path(file_path)
        if path.suffix.lower() not in {".ppt", ".pptx"} or not self.powerpoint_renderer.can_render(path):
            return
        try:
            token = self.file_cache_token(str(path))
            read_preview_cache(str(path), request_page_no(page_no), token, self.powerpoint_renderer.name)
        except OSError:
            return

    def _render_presentation(self, path: Path, page_no: int, token: str, matched_text: str) -> PreviewResult:
        for renderer in (self.preconverted_renderer, self.powerpoint_renderer, self.pptx_text_renderer):
            if not renderer.can_render(path):
                continue
            try:
                return renderer.render(path, page_no, token, matched_text)
            except Exception as exc:
                last_error = exc
        warning = f"프레젠테이션 미리보기를 만들 수 없어 검색 문단을 표시합니다: {last_error}" if "last_error" in locals() else ""
        return PreviewResult("text", matched_text, 1, "matched-text", page_no, warning=warning)


class PdfPageRenderer:
    name = "pdf"

    def can_render(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == ".pdf"

    def render(self, file_path: Path, page_no: int, token: str, matched_text: str = "") -> PreviewResult:
        cached = read_preview_cache(str(file_path), page_no, token, self.name)
        if cached is not None:
            image_bytes, total = cached
            return PreviewResult("image", image_bytes, total, self.name, page_no)

        import fitz

        with fitz.open(file_path) as document:
            if document.needs_pass:
                raise RuntimeError("암호가 걸려 있는 PDF는 미리보기할 수 없습니다.")
            if len(document) == 0:
                raise RuntimeError("미리보기할 페이지가 없습니다.")
            page_index = min(max(page_no - 1, 0), len(document) - 1)
            pixmap = document[page_index].get_pixmap(matrix=fitz.Matrix(1.25, 1.25), alpha=False)
            image_bytes = pixmap.tobytes("png")
            write_preview_cache(str(file_path), page_no, token, self.name, image_bytes, len(document))
            return PreviewResult("image", image_bytes, len(document), self.name, page_no)


class PreconvertedPreviewRenderer:
    name = "preconverted"

    def __init__(self, pdf_renderer: PdfPageRenderer) -> None:
        self.pdf_renderer = pdf_renderer

    def can_render(self, file_path: Path) -> bool:
        if file_path.suffix.lower() not in {".ppt", ".pptx"}:
            return False
        return self._pdf_path(file_path).exists() or self._count_png_pages(file_path) > 0

    def render(self, file_path: Path, page_no: int, token: str, matched_text: str = "") -> PreviewResult:
        png_path = self._png_path(file_path, page_no)
        if png_path is not None:
            total = self._count_png_pages(file_path) or page_no
            return PreviewResult("image", png_path.read_bytes(), total, self.name, page_no)
        pdf_path = self._pdf_path(file_path)
        if pdf_path.exists():
            pdf_token = f"{token}:{pdf_path.stat().st_mtime_ns}:{pdf_path.stat().st_size}"
            result = self.pdf_renderer.render(pdf_path, page_no, pdf_token, matched_text)
            return PreviewResult(result.kind, result.content, result.total_pages, self.name, result.page_no)
        raise FileNotFoundError("preconverted preview not found")

    def _pdf_path(self, file_path: Path) -> Path:
        return file_path.with_suffix(".pdf")

    def _png_path(self, file_path: Path, page_no: int) -> Path | None:
        stem = file_path.stem
        candidates = [
            file_path.with_name(f"{stem}.png") if page_no == 1 else None,
            file_path.with_name(f"{stem}_{page_no}.png"),
            file_path.with_name(f"{stem}_{page_no:03d}.png"),
            file_path.with_name(f"{stem}_page_{page_no}.png"),
            file_path.with_name(f"{stem}_page_{page_no:03d}.png"),
            file_path.with_name(f"{stem}_slide_{page_no}.png"),
            file_path.with_name(f"{stem}_slide_{page_no:03d}.png"),
        ]
        for candidate in candidates:
            if candidate is not None and candidate.exists():
                return candidate
        return None

    def _count_png_pages(self, file_path: Path) -> int:
        stem = file_path.stem
        patterns = [f"{stem}_*.png", f"{stem}_page_*.png", f"{stem}_slide_*.png"]
        matches = set()
        for pattern in patterns:
            matches.update(file_path.parent.glob(pattern))
        if file_path.with_name(f"{stem}.png").exists():
            matches.add(file_path.with_name(f"{stem}.png"))
        return len(matches)


class PowerPointComRenderer:
    name = "powerpoint-com"

    def can_render(self, file_path: Path) -> bool:
        return platform.system().lower() == "windows" and file_path.suffix.lower() in {".ppt", ".pptx"} and self._imports_available()

    def render(self, file_path: Path, page_no: int, token: str, matched_text: str = "") -> PreviewResult:
        cached = read_preview_cache(str(file_path), page_no, token, self.name)
        if cached is not None:
            image_bytes, total = cached
            return PreviewResult("image", image_bytes, total, self.name, page_no)

        import pythoncom
        import win32com.client

        pythoncom.CoInitialize()
        powerpoint = None
        presentation = None
        try:
            powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
            try:
                powerpoint.DisplayAlerts = 0
            except Exception:
                pass
            presentation = powerpoint.Presentations.Open(str(file_path), ReadOnly=True, Untitled=False, WithWindow=False)
            total = int(presentation.Slides.Count)
            if total == 0:
                raise RuntimeError("미리보기할 슬라이드가 없습니다.")
            slide_index = min(max(page_no, 1), total)
            slide_width = float(presentation.PageSetup.SlideWidth or 1280)
            slide_height = float(presentation.PageSetup.SlideHeight or 720)
            export_height = max(1, int(PREVIEW_RENDER_WIDTH * slide_height / max(slide_width, 1)))
            output_path, count_path = preview_cache_paths(str(file_path), page_no, token, self.name)
            presentation.Slides(slide_index).Export(str(output_path), "PNG", PREVIEW_RENDER_WIDTH, export_height)
            count_path.write_text(str(total), encoding="utf-8")
            return PreviewResult("image", output_path.read_bytes(), total, self.name, page_no)
        finally:
            if presentation is not None:
                try:
                    presentation.Close()
                except Exception:
                    pass
            if powerpoint is not None:
                try:
                    powerpoint.Quit()
                except Exception:
                    pass
            pythoncom.CoUninitialize()

    def _imports_available(self) -> bool:
        try:
            import pythoncom  # noqa: F401
            import win32com.client  # noqa: F401
        except Exception:
            return False
        return True


class PptxTextRenderer:
    name = "pptx-text"

    def can_render(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == ".pptx"

    def render(self, file_path: Path, page_no: int, token: str, matched_text: str = "") -> PreviewResult:
        from pptx import Presentation

        presentation = Presentation(file_path)
        total = len(presentation.slides)
        if total == 0:
            return PreviewResult("text", "", 0, self.name, page_no)
        slide = presentation.slides[min(max(page_no - 1, 0), total - 1)]
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        return PreviewResult("text", "\n".join(texts), total, self.name, page_no)


def request_page_no(value: object) -> int:
    try:
        return max(int(value or 1), 1)
    except (TypeError, ValueError):
        return 1


def preview_cache_dir() -> Path:
    output_dir = Path(tempfile.gettempdir()) / "llm_reference_preview"
    output_dir.mkdir(parents=True, exist_ok=True)
    return output_dir


def preview_cache_key(file_path: str, page_no: int, token: str, renderer: str) -> str:
    raw = f"{renderer}:{file_path}:{page_no}:{token}:{PREVIEW_RENDER_WIDTH}"
    return hashlib.sha1(raw.encode("utf-8", errors="ignore")).hexdigest()


def preview_cache_paths(file_path: str, page_no: int, token: str, renderer: str) -> tuple[Path, Path]:
    cache_key = preview_cache_key(file_path, page_no, token, renderer)
    output_dir = preview_cache_dir()
    return output_dir / f"{cache_key}.png", output_dir / f"{cache_key}.count"


def read_preview_cache(file_path: str, page_no: int, token: str, renderer: str) -> tuple[bytes, int] | None:
    output_path, count_path = preview_cache_paths(file_path, page_no, token, renderer)
    if not output_path.exists() or not count_path.exists():
        return None
    try:
        return output_path.read_bytes(), int(count_path.read_text(encoding="utf-8").strip())
    except (OSError, ValueError):
        return None


def write_preview_cache(file_path: str, page_no: int, token: str, renderer: str, image_bytes: bytes, total: int) -> None:
    output_path, count_path = preview_cache_paths(file_path, page_no, token, renderer)
    output_path.write_bytes(image_bytes)
    count_path.write_text(str(total), encoding="utf-8")


def image_bytes_from_pil(image) -> bytes:
    image_buffer = io.BytesIO()
    image.save(image_buffer, format="PNG", optimize=True)
    return image_buffer.getvalue()
