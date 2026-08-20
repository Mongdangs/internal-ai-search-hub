from __future__ import annotations

from pathlib import Path
import zipfile

from src.ingestion.exceptions import PasswordProtectedDocument
from src.ingestion.heading_extractor import enrich_units_with_headings
from src.models import ParsedUnit
from src.utils.page_label import extract_display_page_label
from src.utils.text_cleaner import clean_text


class PptxParser:
    supported = (".pptx",)

    def parse(self, path: str | Path) -> list[ParsedUnit]:
        from pptx import Presentation

        if _is_encrypted_pptx(Path(path)):
            raise PasswordProtectedDocument(f"Password protected PPTX skipped: {path}")
        presentation = Presentation(str(path))
        units: list[ParsedUnit] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = clean_text(shape.text)
                    if text:
                        texts.append(text)
            slide_text = clean_text("\n".join(texts))
            if slide_text:
                section_title = texts[0] if texts else ""
                units.append(
                    ParsedUnit(
                        page_no=index,
                        text=slide_text,
                        section_title=section_title,
                        display_page=extract_display_page_label(slide_text, index),
                    )
                )
        return enrich_units_with_headings(units)


def _is_encrypted_pptx(path: Path) -> bool:
    with path.open("rb") as f:
        header = f.read(8)
    if header.startswith(b"\xd0\xcf\x11\xe0"):
        return True
    try:
        with zipfile.ZipFile(path) as zf:
            names = {name.lower() for name in zf.namelist()}
            return "encryptioninfo" in names or "encryptedpackage" in names
    except zipfile.BadZipFile:
        return False
